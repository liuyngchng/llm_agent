#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# Copyright (c) [2025] [liuyngchng@hotmail.com] - All rights reserved.
import json
import os
import logging.config
import sys
from typing import Generator

import subprocess
import tempfile

import requests

from apps.doc_forge.code_executor import extract_python_blocks, execute_code

import urllib3
# 全局禁用不安全请求警告 InsecureRequestWarning
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

ALLOWED_EXTENSIONS = {
    'txt', 'md', 'csv', 'py', 'js', 'html', 'css', 'json',
    'pdf', 'xlsx', 'xls', 'docx', 'doc', 'ppt', 'pptx',
    'jpg', 'jpeg', 'png', 'gif'
}

# 旧 Office 格式 -> 新格式映射（LibreOffice headless 自动转换）
_OLD_FORMAT_MAP = {'doc': 'docx', 'ppt': 'pptx', 'xls': 'xlsx'}

MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB

log_config_path = 'logging.conf'
if os.path.exists(log_config_path):
    logging.config.fileConfig(log_config_path, encoding="utf-8")
else:
    from common.const import LOG_FORMATTER
    logging.basicConfig(level=logging.INFO,format=LOG_FORMATTER,force=True)
logger = logging.getLogger(__name__)

# ============== 配置常量 ==============
# LLM API 配置（可替换为任何兼容OpenAI API的接口）
class LLMConfig:

    # 请求参数配置
    MAX_TOKENS = int(os.getenv("MAX_TOKENS", 8000))
    TEMPERATURE = float(os.getenv("TEMPERATURE", 0.7))
    TOP_P = float(os.getenv("TOP_P", 0.9))

    # 流式响应配置
    STREAM = True
    TIMEOUT = 60  # 请求超时时间（秒）

    # 系统提示词
    SYSTEM_PROMPT = os.getenv("SYSTEM_PROMPT", "你是一个有用的AI助手。请用中文回答用户的问题。")


def build_doc_processing_system_prompt(file_paths: list[str],
                                        output_dir: str,
                                        upload_dir: str = "upload_doc") -> str:
    """Build system prompt for document processing agent."""
    files_section = ""
    if file_paths:
        files_list = "\n".join(f"  - {fp}" for fp in file_paths)
        files_section = f"""
## 可用文件
以下文件已上传并可供处理（**必须使用下表列出的绝对路径，不要自己拼接路径**）：
{files_list}
脚本执行环境中已预定义了变量 FILE_1、FILE_2...（按顺序对应每个文件的绝对路径），直接在代码中使用。"""

    return f"""你是文档处理助手，帮助用户读取、修改、合并、创建文档（doc/docx、ppt/pptx、xls/xlsx、pdf、csv）。

## 环境与输入
- `UPLOAD_DIR` = `{upload_dir}` — 上传文件所在目录
- `OUTPUT_DIR` = `{output_dir}` — 所有输出文件**必须**保存到此目录
- `FILE_1`, `FILE_2`, ... — 预定义变量，指向每个上传文件的**绝对路径**，直接使用即可
- 旧格式（.doc/.ppt/.xls）上传时已自动转为新格式（.docx/.pptx/.xlsx）

## 可用库

**Word (docx):**
- `python-docx` — 读取、修改、创建 Word 文档。可操作段落、表格、样式、图片、页眉页脚
- `docxtpl` — 基于 Jinja2 模板填充 Word 文档

**PowerPoint (pptx):**
- `python-pptx` — 读取、修改、创建 PPT 演示文稿。可操作幻灯片、形状、表格、图片

**Excel (xlsx):**
- `openpyxl` — 读取、修改、创建 Excel 工作簿。可操作工作表、单元格、公式、样式、图表

**PDF:**
- `pdfplumber` — 提取 PDF 中的文本和表格
- `pypdf` — 修改已有 PDF：合并、拆分、提取页面、旋转、加水印
- `reportlab` — 以编程方式创建新 PDF

**其他:**
- `Pillow (PIL)` — 图片处理
- `pandas` — 数据处理和分析
- `csv` — CSV 文件读写
- `common.ocr_util.ImageOCR` — 图片 OCR 文字识别

## 强制规则
1. **直接操作原生格式**：修改文档必须使用对应格式的原生库打开文件、修改内容、保存到 `OUTPUT_DIR`
2. **禁止 markdown 中转**：严禁先将文档转成 markdown 再转回，这会导致格式丢失
3. **输出到指定目录**：所有生成的文档必须保存到 `OUTPUT_DIR`（即 `{output_dir}`）。使用 `os.path.join(OUTPUT_DIR, "文件名.扩展名")` 构造输出路径
4. **一个代码块输出**：将所有操作写在一个 ```python 代码块中，第一行用注释简要说明功能。代码块是触发文档修改的唯一方式
5. **用 print() 汇报结果**：脚本执行完成后用 `print()` 输出操作结果和生成的文件名{files_section}"""


def allowed_file(filename):
    """检查文件扩展名是否允许"""
    return '.' in filename and \
        filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def _convert_old_format(filepath: str) -> str:
    """使用 LibreOffice headless 将旧 Office 格式转为新格式，返回转换后文件路径。"""
    ext = os.path.splitext(filepath)[1].lstrip('.').lower()
    target_ext = _OLD_FORMAT_MAP[ext]
    base_name = os.path.splitext(os.path.basename(filepath))[0]

    output_dir = tempfile.mkdtemp(prefix='docforge_')
    result = subprocess.run(
        ['libreoffice', '--headless', '--convert-to', target_ext,
         '--outdir', output_dir, filepath],
        capture_output=True, text=True, timeout=60,
    )
    if result.returncode != 0:
        error_msg = result.stderr.strip() or result.stdout.strip()
        logger.error(f"LibreOffice 转换失败: {filepath} -> {target_ext}: {error_msg}")
        raise RuntimeError(f"旧格式文件转换失败（{ext} -> {target_ext}）")

    output_path = os.path.join(output_dir, f"{base_name}.{target_ext}")
    logger.info(f"旧格式转换成功: {os.path.basename(filepath)} -> {os.path.basename(output_path)}")
    return output_path


def get_pptx_content(filepath) -> str:
    """使用 python-pptx 原生接口提取 pptx 文件内容。"""
    try:
        from pptx import Presentation

        prs = Presentation(filepath)
        slides_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_lines = [f"## Slide {slide_num}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_lines.append(shape.text.strip())
                if shape.has_table:
                    table = shape.table
                    rows_text = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows_text.append(" | ".join(cells))
                    if rows_text:
                        slide_lines.append("\n".join(rows_text))
            slides_text.append("\n".join(slide_lines))
        result = "\n\n".join(slides_text)
        logger.info(f"pptx_file_parse_success, {filepath}，slides: {len(prs.slides)}，len: {len(result)}")
        return result

    except ImportError:
        logger.warning("需要安装 python-pptx 库来解析PPT文件")
        return "[需要安装 python-pptx 库来解析PPT文件]"
    except Exception as e:
        logger.error(f"读取pptx文件时出错: {str(e)}")
        return f"[读取pptx文件时出错: {str(e)}]"


def extract_text_from_file(filepath, filename, ocr_engine=None):
    """从文件中提取文本内容"""
    try:
        ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''

        logger.info(f"正在提取文件: {filename}, 扩展名: {ext}")

        # 文本文件直接读取
        if ext in ['txt', 'md', 'py', 'js', 'html', 'css', 'json']:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                logger.info(f"文本文件读取成功，长度: {len(content)} 字符")
                return content

        # CSV 文件 — 解析为 Markdown 表格
        elif ext == 'csv':
            try:
                import csv as _csv
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    reader = _csv.reader(f)
                    rows = list(reader)
                if not rows:
                    return "[CSV 文件为空]"
                lines = []
                lines.append('| ' + ' | '.join(rows[0]) + ' |')
                lines.append('| ' + ' | '.join(['---'] * len(rows[0])) + ' |')
                for row in rows[1:]:
                    # 补齐缺列
                    padded = row + [''] * (len(rows[0]) - len(row))
                    lines.append('| ' + ' | '.join(padded[:len(rows[0])]) + ' |')
                result = '\n'.join(lines)
                logger.info(f"CSV文件解析成功，行数: {len(rows)}")
                return result
            except Exception as e:
                logger.error(f"CSV解析失败: {str(e)}")
                return f"[CSV 解析失败: {str(e)}]"

        # PDF文件（需要安装PyPDF2或pdfplumber）
        elif ext == 'pdf':
            try:
                import pdfplumber
                text = []
                with pdfplumber.open(filepath) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:  # 添加判断，避免None值
                            text.append(page_text)
                result = '\n'.join(text)
                logger.info(f"pdf_file_parse_success, {filepath}，page_count: {len(pdf.pages)}，file_len: {len(result)}")
                return result
            except ImportError:
                logger.warning("需要安装 pdfplumber 库来解析PDF文件")
                return "[需要安装 pdfplumber 库来解析PDF文件]"

        # Word 文档（.doc 自动通过 LibreOffice 转为 .docx）
        elif ext in ('docx', 'doc'):
            path = _convert_old_format(filepath) if ext == 'doc' else filepath
            return get_docx_content(path)

        # PPT 文档（.ppt 自动通过 LibreOffice 转为 .pptx）
        elif ext in ('pptx', 'ppt'):
            path = _convert_old_format(filepath) if ext == 'ppt' else filepath
            return get_pptx_content(path)

        # Excel 文件（.xls 自动通过 LibreOffice 转为 .xlsx）
        elif ext in ('xlsx', 'xls'):
            path = _convert_old_format(filepath) if ext == 'xls' else filepath
            return get_xlsx_content(path)

        # 图片文件 — 使用 OCR API 识别文字
        elif ext in ['jpg', 'jpeg', 'png', 'gif']:
            if ocr_engine:
                try:
                    result = ocr_engine.extract_text_from_image(filepath)
                    if result['success']:
                        text = result['text']
                        logger.info(f"OCR API 识别成功，识别到 {len(text)} 字符")
                        return text
                    else:
                        logger.warning(f"OCR API 识别失败: {result.get('error')}")
                        return f"[OCR 识别失败: {result.get('error')}]"
                except Exception as e:
                    logger.error(f"OCR API 调用异常: {str(e)}")
                    return f"[OCR 识别异常: {str(e)}]"
            else:
                return "[未配置 OCR 引擎，无法识别图片文字]"

        else:
            logger.warning(f"不支持的文件格式: {ext}")
            return f"[不支持的文件格式: {ext}]"

    except Exception as e:
        logger.error(f"读取文件时出错: {str(e)}")
        return f"[读取文件时出错: {str(e)}]"


def get_xlsx_content(filepath) -> str:
    """
    获取 xlsx 文件的内容
    """
    try:
        from openpyxl import load_workbook
        wb = load_workbook(filename=filepath, read_only=True, data_only=True)
        text_lines = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_lines.append(f"## {sheet_name}")

            # 获取最大列宽
            max_col = ws.max_column
            max_row = ws.max_row

            if max_row > 0 and max_col > 0:
                # 创建Markdown表格 - 修复这里
                rows = []
                for row in ws.iter_rows(min_row=1, max_row=max_row,
                                        min_col=1, max_col=max_col,
                                        values_only=True):
                    # 不要转义管道符，而是处理空值和特殊字符
                    formatted_row = []
                    for cell in row:
                        if cell is None:
                            formatted_row.append("")
                        else:
                            # 将单元格内容转换为字符串，并处理换行
                            cell_str = str(cell)
                            # 替换换行符为空格，避免破坏表格格式
                            cell_str = cell_str.replace('\n', ' ').replace('\r', ' ')
                            # 清理多余的空白字符
                            cell_str = ' '.join(cell_str.split())
                            formatted_row.append(cell_str)
                    rows.append(formatted_row)

                # 生成Markdown表格
                if rows:
                    # 表头
                    md_lines = ['| ' + ' | '.join(rows[0]) + ' |']
                    # 表头分隔线
                    header_separator = ['---'] * len(rows[0])
                    md_lines.append('| ' + ' | '.join(header_separator) + ' |')
                    # 数据行
                    for row in rows[1:]:
                        md_lines.append('| ' + ' | '.join(row) + ' |')

                    text_lines.append('\n'.join(md_lines))
                text_lines.append('')  # 空行分隔

        result = '\n'.join(text_lines)
        logger.info(f"Excel文件解析成功，工作表数: {len(wb.sheetnames)}，长度: {len(result)} 字符")
        return result
    except ImportError:
        logger.warning("需要安装openpyxl库来解析Excel文件")
        return "[需要安装openpyxl库来解析Excel文件]"
    except Exception as e:
        logger.error(f"读取Excel文件时出错: {str(e)}")
        return f"[读取Excel文件时出错: {str(e)}]"


def get_docx_content(filepath) -> str:
    """
    使用 python-docx 原生接口提取 docx 文件内容。
    保留文档结构（标题层级、表格、列表），不通过 markdown 转换。
    """
    try:
        from docx import Document

        doc = Document(filepath)
        text_parts = []

        for para in doc.paragraphs:
            style_name = para.style.name if para.style else ''
            text = para.text.strip()
            if not text:
                continue

            if style_name.startswith('Heading'):
                try:
                    level = int(style_name.replace('Heading ', ''))
                    prefix = '#' * min(level, 6)
                    text_parts.append(f"\n{prefix} {text}")
                except ValueError:
                    text_parts.append(f"\n## {text}")
            elif 'List' in style_name or 'list' in style_name.lower():
                text_parts.append(f"- {text}")
            else:
                text_parts.append(text)

        for i, table in enumerate(doc.tables):
            text_parts.append(f"\n### 表格 {i + 1}")
            row_data = []
            for row in table.rows:
                cells = [cell.text.strip().replace('\n', ' ').replace('|', '\\|') for cell in row.cells]
                row_data.append('| ' + ' | '.join(cells) + ' |')
            if row_data:
                text_parts.append(row_data[0])
                num_cols = len(table.rows[0].cells)
                text_parts.append('| ' + ' | '.join(['---'] * num_cols) + ' |')
                for row in row_data[1:]:
                    text_parts.append(row)

        for section in doc.sections:
            header = section.header
            if header and header.paragraphs:
                h_text = ' '.join(p.text for p in header.paragraphs if p.text.strip())
                if h_text:
                    text_parts.insert(0, f"[页眉: {h_text}]")
            footer = section.footer
            if footer and footer.paragraphs:
                f_text = ' '.join(p.text for p in footer.paragraphs if p.text.strip())
                if f_text:
                    text_parts.append(f"[页脚: {f_text}]")

        result = '\n'.join(text_parts)
        logger.info(f"docx_file_parse_success, {filepath}, file_len: {len(result)}")
        return result

    except ImportError:
        logger.warning("需要安装 python-docx 库来解析Word文档")
        return "[需要安装 python-docx 库来解析Word文档]"
    except Exception as e:
        logger.error(f"读取docx文件时出错: {str(e)}")
        return f"[读取docx文件时出错: {str(e)}]"


def generate_stream_response(messages: list, llm_cfg: dict,
                            include_done: bool = True) -> Generator[str, None, None]:
    """
    生成流式响应
    :param include_done: 是否在流结束后发送 [DONE] 信号
    """
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {llm_cfg['llm_api_key']}"
    }

    payload = {
        "model": llm_cfg['llm_model_name'],
        "messages": messages,
        "max_tokens": LLMConfig.MAX_TOKENS,
        "temperature": LLMConfig.TEMPERATURE,
        "top_p": LLMConfig.TOP_P,
        "stream": LLMConfig.STREAM,
    }

    try:
        logger.info(f"向LLM API发送请求，模型: {llm_cfg['llm_model_name']}, 消息数量: {len(messages)}")
        logger.info(f"请求参数: max_tokens={payload['max_tokens']}, temperature={LLMConfig.TEMPERATURE}")

        response = requests.post(
            f"{llm_cfg['llm_api_uri']}/chat/completions",
            headers=headers,
            json=payload,
            stream=True,
            verify=False,
            timeout=LLMConfig.TIMEOUT
        )

        # 记录响应状态码
        logger.info(f"llm_api_response_code: {response.status_code}")

        # 如果响应不是200，记录详细错误信息
        if response.status_code != 200:
            error_content = response.text
            logger.error(f"llm_api_response_error, {response.status_code}: {error_content}")

            try:
                error_json = response.json()
                logger.error(f"llm_api_response_error_json: {json.dumps(error_json, indent=2, ensure_ascii=False)}")

                # 提取具体错误信息
                if 'error' in error_json:
                    error_msg = error_json['error']
                    if isinstance(error_msg, dict) and 'message' in error_msg:
                        error_detail = error_msg['message']
                    elif isinstance(error_msg, str):
                        error_detail = error_msg
                    else:
                        error_detail = str(error_msg)

                    logger.error(f"错误消息: {error_detail}")
                else:
                    logger.error(f"完整错误响应: {json.dumps(error_json, ensure_ascii=False)}")

            except json.JSONDecodeError:
                logger.error(f"错误响应不是JSON格式，原始内容: {error_content}")
            except Exception as e:
                logger.error(f"解析错误响应时出错: {str(e)}, 原始内容: {error_content}")

            # 向客户端返回错误信息
            error_message = f"API请求失败: {response.status_code} - {response.reason}"
            yield f"data: {json.dumps({'error': error_message})}\n\n"
            if include_done:
                yield "data: [DONE]\n\n"
            return

        response.raise_for_status()

        for line in response.iter_lines():
            if line:
                line = line.decode('utf-8')
                if line.startswith('data: '):
                    data = line[6:]  # 移除 'data: ' 前缀
                    if data != '[DONE]':
                        try:
                            chunk = json.loads(data)
                            if 'choices' in chunk and len(chunk['choices']) > 0:
                                delta = chunk['choices'][0].get('delta', {})
                                if 'content' in delta and delta['content']:
                                    yield f"data: {json.dumps({'content': delta['content']})}\n\n"
                        except json.JSONDecodeError:
                            logger.warning("解析JSON数据失败")
                            continue

        if include_done:
            yield "data: [DONE]\n\n"

    except requests.exceptions.RequestException as e:
        error_msg = f"API请求错误: {str(e)}"
        logger.error(error_msg)

        # 如果有响应对象，记录更多信息
        if hasattr(e, 'response') and e.response is not None:
            try:
                error_json = e.response.json()
                logger.error(f"请求异常中的错误详情: {json.dumps(error_json, ensure_ascii=False)}")
            except:
                logger.error(f"请求异常中的响应内容: {e.response.text}")

        yield f"data: {json.dumps({'error': error_msg})}\n\n"
        if include_done:
            yield "data: [DONE]\n\n"


def _call_llm_sync(messages: list, llm_cfg: dict) -> str:
    """同步调用 LLM，收集完整响应文本（非流式）。"""
    full_text = ""
    for chunk in generate_stream_response(messages, llm_cfg, include_done=False):
        if chunk.startswith('data: ') and chunk != 'data: [DONE]\n\n':
            try:
                data = json.loads(chunk[6:].strip())
                if 'content' in data:
                    full_text += data['content']
            except (json.JSONDecodeError, KeyError):
                pass
    return full_text


def generate_stream_response_with_execution(
    messages: list, llm_cfg: dict, output_dir: str = "output_doc",
    upload_dir: str = "upload_doc", file_paths: list | None = None
) -> Generator[str, None, None]:
    """
    流式响应 + 代码执行引擎。

    先流式传输 LLM 回复，然后从回复中提取 Python 代码块并在服务器端执行。
    执行结果（stdout/stderr/保存信息）作为额外的 SSE 数据发送给前端。
    """
    import re as _re

    full_response = ""

    # Phase 1: 收集 LLM 完整回复（不直接流式展示，以便过滤代码块）
    for sse_chunk in generate_stream_response(messages, llm_cfg, include_done=False):
        if sse_chunk.startswith('data: ') and sse_chunk != 'data: [DONE]\n\n':
            try:
                chunk_data = json.loads(sse_chunk[6:].strip())
                if 'error' in chunk_data:
                    # API 错误直接传给前端，不继续执行
                    yield sse_chunk
                    yield "data: [DONE]\n\n"
                    return
                if 'content' in chunk_data:
                    full_response += chunk_data['content']
            except (json.JSONDecodeError, KeyError):
                pass

    # 将代码块替换为可折叠详情块，同时收集代码用于后续执行（只提取一次，避免二次正则不一致）
    collected_codes = []
    def _summarize_code_block(match: _re.Match) -> str:
        code = match.group(1).strip()
        collected_codes.append(code)  # 边替换边收集
        line_count = code.count('\n') + 1
        first_line = code.split('\n')[0] if code else ''
        desc = ""
        if first_line.lstrip().startswith('#'):
            desc = first_line.lstrip('#').strip()
        summary = desc if desc else f"已生成处理脚本（{line_count} 行），正在执行..."
        escaped_code = code.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        return f"""
<details class="script-details">
<summary>📝 <b>脚本</b>（{line_count} 行）：{summary}</summary>

<pre><code class="language-python">{escaped_code}</code></pre>
</details>
"""

    # 只截取最后一个代码块之前的文本作为"方案"，不展示LLM的虚假完成声明
    plan_text = full_response
    last_match = None
    for m in _re.finditer(r'```python\s*\n.*?\n\s*```', full_response, _re.DOTALL):
        last_match = m
    if last_match:
        plan_text = full_response[:last_match.end()]

    display_response = _re.sub(
        r'```python\s*\n(.*?)\n\s*```',
        _summarize_code_block,
        plan_text,
        flags=_re.DOTALL
    )

    yield f"data: {json.dumps({'content': display_response})}\n\n"

    # Phase 2: 执行脚本（合并所有代码块为一个脚本，避免多进程变量不共享）
    MAX_RETRIES = 2
    code_blocks = collected_codes
    if code_blocks:
        logger.info(f"从LLM回复中提取到 {len(code_blocks)} 个Python代码块，合并执行...")
        combined_code = "\n\n".join(code_blocks)
        retry_history = messages + [{"role": "assistant", "content": full_response}]

        result = None
        for attempt in range(MAX_RETRIES):
            if attempt > 0:
                logger.info(f"脚本执行出错，第{attempt}次重试...")
            result = execute_code(combined_code, output_dir=output_dir, upload_dir=upload_dir, file_paths=file_paths)

            if result['success'] and not result['stderr']:
                break

            stderr = result.get('stderr', '')
            # 缺少模块时不重试
            if _re.search(r"ModuleNotFoundError: No module named '([^']+)'", stderr):
                logger.warning(f"缺少模块，无法自动修复: {stderr}")
                break

            if attempt < MAX_RETRIES - 1:
                error_info = result['stderr'] or result['stdout'] or f"返回码: {result['returncode']}"
                retry_msg = (
                    f"脚本执行出错，请根据错误信息修复代码：\n\n"
                    f"```\n{error_info}\n```\n\n"
                    f"请在 ```python 代码块中输出修复后的完整脚本。"
                )
                retry_history.append({"role": "user", "content": retry_msg})
                fix_response = _call_llm_sync(retry_history, llm_cfg)
                fixed_blocks = extract_python_blocks(fix_response)
                if fixed_blocks:
                    retry_history.append({"role": "assistant", "content": fix_response})
                    combined_code = "\n\n".join(fixed_blocks)
                else:
                    logger.info("LLM 修复回复中未包含代码块，停止重试")
                    break

        # 汇总结果
        output_parts = []
        if result and result['success']:
            new_files = result.get('new_files', [])
            if new_files:
                file_names = ", ".join(sorted(new_files))
                output_parts.append(f"✅ 修改完成！处理后的文件：`{file_names}`")
                if result['stdout']:
                    output_parts.append(f"**处理日志:**\n```\n{result['stdout'][-500:]}\n```")
            else:
                output_parts.append(
                    "📋 脚本已执行完毕，但未生成新的文档文件。\n\n"
                    "请补充你的修改要求，比如：\n"
                    "- 需要修改哪些内容？\n"
                    "- 需要调整哪些格式？\n"
                    "- 需要输出什么格式的文件？"
                )
        else:
            output_parts.append("❌ 处理失败，请根据错误信息调整你的要求后重试。")
            if result and result.get('stderr'):
                output_parts.append(f"**错误信息:**\n```\n{result['stderr'][-1000:]}\n```")

        yield f"data: {json.dumps({'content': '\n\n'.join(output_parts) + '\n\n'})}\n\n"

    else:
        logger.info("LLM回复中未提取到Python代码块，跳过代码执行")
        fake_keywords = ['修改完成', '已修改', '脚本执行', '文件已保存', '文档已', '已保存', '已修正']
        if any(kw in full_response for kw in fake_keywords):
            logger.warning("LLM 声称已修改但未输出代码块，尝试一次修复")
            retry_history = messages + [{"role": "assistant", "content": full_response}]
            retry_msg = (
                "你声称已修改文档，但没有输出 Python 代码块。"
                "请在 ```python 代码块中输出可执行的脚本来完成修改。"
            )
            retry_history.append({"role": "user", "content": retry_msg})
            fix_response = _call_llm_sync(retry_history, llm_cfg)
            fixed_blocks = extract_python_blocks(fix_response)
            if fixed_blocks:
                result = execute_code("\n\n".join(fixed_blocks), output_dir=output_dir, upload_dir=upload_dir, file_paths=file_paths)
                output_parts = []
                if result['success'] and result['new_files']:
                    file_names = ", ".join(result['new_files'])
                    output_parts.append(f"✅ 修改完成！处理后的文件：`{file_names}`")
                else:
                    output_parts.append("❌ 脚本执行未成功，请重新描述你的修改需求。")
                    if result.get('stderr'):
                        output_parts.append(f"**错误信息:**\n```\n{result['stderr'][-500:]}\n```")
                if result.get('stdout'):
                    output_parts.append(f"**输出:**\n```\n{result['stdout'][-500:]}\n```")
                yield f"data: {json.dumps({'content': '\n\n'.join(output_parts) + '\n\n'})}\n\n"
            else:
                warning = (
                    "\n\n---\n\n⚠️ 未能生成可执行的代码。请重新发送消息，明确要求 AI 输出 Python 代码块。\n\n"
                )
                logger.warning("一次修复仍未生成代码块，放弃")
                yield f"data: {json.dumps({'content': warning})}\n\n"

    yield "data: [DONE]\n\n"
